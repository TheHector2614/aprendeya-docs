import re
from pathlib import Path

from bs4 import BeautifulSoup
import markdown
from docx import Document
from openpyxl import load_workbook


class Extractor:
    def __init__(self, ocr_fallback: bool = True):
        self.ocr_fallback = ocr_fallback

    def extract(self, path: Path) -> str:
        ext = path.suffix.lower()
        handler = {
            ".html": self._extract_html,
            ".htm": self._extract_html,
            ".md": self._extract_markdown,
            ".docx": self._extract_docx,
            ".xlsx": self._extract_xlsx,
            ".pptx": self._extract_pptx,
            ".pdf": self._extract_pdf,
            ".txt": self._extract_txt,
            ".csv": self._extract_csv,
            ".json": self._extract_json,
        }
        handler_fn = handler.get(ext)
        if not handler_fn:
            raise ValueError(f"Formato no soportado: {ext}")
        return handler_fn(path)

    def _extract_html(self, path: Path) -> str:
        soup = BeautifulSoup(path.read_text("utf-8"), "lxml")
        for tag in soup(["script", "style", "nav", "header", "footer", "noscript"]):
            tag.decompose()
        lines = []
        for el in soup.body.descendants if soup.body else soup.children:
            if el.name in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = int(el.name[1])
                lines.append(f"{'#' * level} {el.get_text(strip=True)}")
            elif el.name == "p":
                t = el.get_text(strip=True)
                if t:
                    lines.append(t)
            elif el.name == "li":
                t = el.get_text(strip=True)
                if t:
                    lines.append(f"- {t}")
        return "\n".join(lines) if lines else "\n".join(
            line.strip() for line in soup.get_text().splitlines() if line.strip()
        )

    def _extract_markdown(self, path: Path) -> str:
        html = markdown.markdown(path.read_text("utf-8"))
        return self._extract_html_from_string(html)

    def _extract_docx(self, path: Path) -> str:
        doc = Document(str(path))
        lines = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                if p.style and p.style.name and "heading" in p.style.name.lower():
                    level = p.style.name.lower().replace("heading ", "")
                    prefix = "#" * min(int(level), 6) if level.isdigit() else "##"
                    lines.append(f"{prefix} {t}")
                else:
                    lines.append(t)
        for table in doc.tables:
            lines.append("\n[tabla]")
            for row in table.rows:
                cells = " | ".join(c.text.strip() for c in row.cells)
                if cells.strip():
                    lines.append(cells)
        return "\n".join(lines)

    def _extract_xlsx(self, path: Path) -> str:
        wb = load_workbook(path, data_only=True)
        lines = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            lines.append(f"\n## {sheet}")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    lines.append(" | ".join(vals))
        return "\n".join(lines)

    def _extract_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise ValueError("python-pptx no instalado. Ejecuta: pip install python-pptx")
        prs = Presentation(str(path))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n--- Diapositiva {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = " | ".join(c.text.strip() for c in row.cells)
                        if cells.strip():
                            lines.append(cells)
            # Notas del orador
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lines.append(f"[notas] {notes}")
        return "\n".join(lines)

    def _extract_pdf(self, path: Path) -> str:
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ValueError("PyMuPDF no instalado. Ejecuta: pip install pymupdf")

        doc = fitz.open(str(path))
        pages = []
        for i, page in enumerate(doc, 1):
            text = page.get_text().strip()
            if text:
                pages.append(f"\n--- Página {i} ---\n{text}")
            elif self.ocr_fallback:
                ocr_text = self._ocr_page(page)
                pages.append(f"\n--- Página {i} (OCR) ---\n{ocr_text}")
        doc.close()
        return "\n".join(pages)

    def _ocr_page(self, page) -> str:
        try:
            import pytesseract
            from PIL import Image
            import io
            pix = page.get_pixmap(dpi=200)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            return pytesseract.image_to_string(img, lang="spa+eng")
        except ImportError:
            return "[OCR no disponible: instala pytesseract + tesseract-ocr]"

    def _extract_txt(self, path: Path) -> str:
        return path.read_text("utf-8")

    def _extract_csv(self, path: Path) -> str:
        import csv
        import io
        text = path.read_text("utf-8-sig")
        reader = csv.reader(io.StringIO(text))
        lines = []
        for i, row in enumerate(reader):
            if any(c.strip() for c in row):
                lines.append(" | ".join(c.strip() for c in row))
        return "\n".join(lines)

    def _extract_json(self, path: Path) -> str:
        import json
        data = json.loads(path.read_text("utf-8"))
        return self._json_to_text(data)

    def _json_to_text(self, data, prefix: str = "") -> str:
        if isinstance(data, dict):
            parts = []
            for k, v in data.items():
                if isinstance(v, (dict, list)):
                    parts.append(self._json_to_text(v, prefix=f"{prefix}{k}."))
                else:
                    parts.append(f"{prefix}{k}: {v}")
            return "\n".join(parts)
        elif isinstance(data, list):
            parts = []
            for i, item in enumerate(data):
                if isinstance(item, (dict, list)):
                    parts.append(self._json_to_text(item, prefix=f"{prefix}[{i}]."))
                else:
                    parts.append(f"{prefix}[{i}]: {item}")
            return "\n".join(parts)
        else:
            return f"{prefix}: {data}"

    def _extract_html_from_string(self, html: str) -> str:
        soup = BeautifulSoup(html, "lxml")
        lines = []
        for el in soup.descendants:
            if el.name in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = int(el.name[1])
                lines.append(f"{'#' * level} {el.get_text(strip=True)}")
            elif el.name == "p":
                t = el.get_text(strip=True)
                if t:
                    lines.append(t)
            elif el.name == "li":
                t = el.get_text(strip=True)
                if t:
                    lines.append(f"- {t}")
        return "\n".join(lines) if lines else "\n".join(
            line.strip() for line in soup.get_text().splitlines() if line.strip()
        )
