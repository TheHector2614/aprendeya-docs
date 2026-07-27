import docx, openpyxl, pptx
from pathlib import Path
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

BASE = Path("raw")

# ── OPE-001.docx ──
d = docx.Document()
d.add_heading("Procedimiento de Atencion al Estudiante", 0)
d.add_heading("Canal de Contacto", 1)
d.add_paragraph("El estudiante puede contactar a traves de los siguientes canales: correo electronico (soporte@aprendeya.com), chat en vivo (disponible 7am-10pm), telefono (01-800-123-4567, horario laboral) y formulario web.")
d.add_heading("Niveles de Atencion", 1)
d.add_paragraph("Nivel 1 (Soporte): Resolucion de consultas basicas. Tiempo maximo de respuesta: 2 horas habiles.")
d.add_paragraph("Nivel 2 (Especializado): Problemas complejos escalados por Nivel 1. Tiempo maximo de respuesta: 24 horas habiles.")
d.add_paragraph("Nivel 3 (Direccion): Incidentes criticos o reclamos no resueltos. Tiempo maximo de respuesta: 48 horas habiles.")
d.add_heading("SLAs", 1)
d.add_paragraph("Consulta general: 2 horas habiles. Problema tecnico: 4 horas habiles. Reclamo o queja: 24 horas habiles. Incidente critico: 15 minutos.")
d.add_heading("Escalamiento", 1)
d.add_paragraph("Si el estudiante no queda satisfecho con la respuesta de Nivel 1, puede solicitar escalamiento a Nivel 2. El supervisor de turno evalua el caso y asigna un especialista en un maximo de 1 hora.")
d.save(str(BASE / "sharepoint" / "ope" / "OPE-001.docx"))
print("[creado] OPE-001.docx")

# ── OPE-003.docx ──
d = docx.Document()
d.add_heading("Manual de Procesos Operativos", 0)
d.add_heading("Gestion de Usuarios", 1)
d.add_paragraph("La creacion de cuentas se automatiza via Auth Service. La eliminacion de cuentas se procesa en un maximo de 48 horas tras la solicitud.")
d.add_heading("Procesamiento de Pagos", 1)
d.add_paragraph("Los pagos se procesan via Payment Service integrado con Stripe y PayPal. Los reembolsos se procesan en un maximo de 5 dias habiles.")
d.add_heading("Publicacion de Cursos", 1)
d.add_paragraph("El instructor crea el contenido en la herramienta de autor. El equipo de calidad revisa y aprueba. El proceso completo toma en promedio 10 dias habiles.")
d.add_heading("Gestion de Certificados", 1)
d.add_paragraph("Los certificados se generan automaticamente al completar el curso con nota minima de 60/100. La generacion toma menos de 5 minutos.")
d.save(str(BASE / "sharepoint" / "ope" / "OPE-003.docx"))
print("[creado] OPE-003.docx")

# ── CAL-001.docx ──
d = docx.Document()
d.add_heading("Procedimiento de Evaluacion de Cursos", 0)
d.add_heading("Indicadores de Calidad", 1)
d.add_paragraph("Los cursos se evaluan trimestralmente: tasa de completitud (meta > 65%), calificacion promedio (meta > 4.0/5.0), NPS (meta > 50), tiempo de respuesta del instructor (meta < 24h).")
d.add_heading("Encuesta de Satisfaccion", 1)
d.add_paragraph("Al finalizar cada curso, el estudiante recibe una encuesta de 8 preguntas anonima que toma menos de 5 minutos.")
d.add_heading("Revision de Contenido", 1)
d.add_paragraph("Cada curso se revisa anualmente. Los cursos con calificacion menor a 3.5/5.0 entran en revision inmediata. El equipo academico tiene 30 dias para presentar un plan de mejora.")
d.add_heading("Acciones Correctivas", 1)
d.add_paragraph("Si un curso recibe mas de 3 quejas en un mes, se suspende temporalmente y se asigna un revisor pedagogico. El instructor recibe retroalimentacion y tiene 15 dias para implementar mejoras.")
d.save(str(BASE / "sharepoint" / "cal" / "CAL-001.docx"))
print("[creado] CAL-001.docx")

# ── RH-003.docx ──
d = docx.Document()
d.add_heading("Politica de Vacaciones y Ausencias", 0)
d.add_heading("Periodo de Vacaciones", 1)
d.add_paragraph("Los colaboradores tienen derecho a 15 dias habiles de vacaciones anuales, acumulables hasta por 2 anos. Deben solicitarse con al menos 15 dias de anticipacion.")
d.add_heading("Ausencias Justificadas", 1)
d.add_paragraph("Se consideran ausencias justificadas: incapacidad medica, calamidad domestica (hasta 3 dias), licencia de maternidad (18 semanas), paternidad (2 semanas), citas medicas (hasta 4 horas).")
d.add_heading("Permisos Remunerados", 1)
d.add_paragraph("Matrimonio (5 dias), nacimiento de hijo (2 dias), fallecimiento de familiar en primer grado (3 dias), tramites de vivienda (1 dia al mes).")
d.add_heading("Registro y Aprobacion", 1)
d.add_paragraph("Todas las ausencias se registran en el sistema de gestion humana. Solicitudes de menos de 2 dias las aprueba el lider directo. Mas de 2 dias requieren aprobacion de RH.")
d.save(str(BASE / "sharepoint" / "rh" / "RH-003.docx"))
print("[creado] RH-003.docx")

# ── RH-004.docx ──
d = docx.Document()
d.add_heading("Politica de Nomina y Compensaciones", 0)
d.add_heading("Estructura Salarial", 1)
d.add_paragraph("Salario base (segun rol y experiencia), bonificacion por desempeno (hasta 20% del salario anual), y beneficios extralegales. Los salarios se revisan anualmente en enero.")
d.add_heading("Frecuencia de Pago", 1)
d.add_paragraph("Nomina quincenal: dias 15 y 30 de cada mes. Pagos via transferencia electronica.")
d.add_heading("Bonificaciones", 1)
d.add_paragraph("Desempeno: semestral (julio y diciembre). Resultados: anual en marzo. Referidos: 1 SMLDV por cada referido que complete el periodo de prueba.")
d.add_heading("Deducciones", 1)
d.add_paragraph("Deducciones de ley: salud (4%), pension (4%), ARL, retencion en la fuente. Deducciones voluntarias autorizadas por el colaborador.")
d.save(str(BASE / "sharepoint" / "rh" / "RH-004.docx"))
print("[creado] RH-004.docx")

# ── FIN-001.xlsx ──
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Facturacion"
ws.append(["Concepto", "Detalle", "Valor"])
ws.append(["Periodo de facturacion", "Mensual / Semestral / Anual", ""])
ws.append(["Dias de gracia", "5 dias calendario despues del vencimiento", ""])
ws.append(["Interes por mora", "1.5% mensual sobre el saldo pendiente", ""])
ws.append(["Medios de pago", "TC, TD, transferencia, PayPal, PSE, Nequi", ""])
ws.append(["Descuento pronto pago", "5% si se paga dentro de los primeros 3 dias", ""])
ws.append(["Plan Basico Mensual", "Acceso a 5 cursos por mes", 49000])
ws.append(["Plan Premium Mensual", "Acceso ilimitado a todo el catalogo", 89000])
ws.append(["Plan Premium Anual", "Acceso ilimitado + certificados gratis", 799000])
ws2 = wb.create_sheet("Cobranzas")
ws2.append(["Etapa", "Tiempo", "Accion"])
ws2.append(["Pre-vencimiento", "3 dias antes", "Recordatorio automatico email y SMS"])
ws2.append(["Vencimiento", "Dia 0", "Notificacion de pago pendiente"])
ws2.append(["Mora temprana", "Dias 1-5", "Recordatorio diario automatico"])
ws2.append(["Mora media", "Dias 6-15", "Llamada de cobranza + email"])
ws2.append(["Mora alta", "Dias 16-30", "Suspension de acceso a cursos nuevos"])
ws2.append(["Mora critica", "Dias 31+", "Suspension total + reporte a centrales"])
wb.save(str(BASE / "drive" / "fin" / "FIN-001.xlsx"))
print("[creado] FIN-001.xlsx")

# ── FIN-002.xlsx ──
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Presupuesto 2026"
ws.append(["Categoria", "Presupuesto Anual", "Ejecutado Q1", "Ejecutado Q2", "% Ejecucion"])
data = [
    ["Tecnologia e Infraestructura", 480000000, 115000000, 120000000, 49],
    ["Desarrollo de Contenido", 350000000, 80000000, 90000000, 49],
    ["Marketing y Publicidad", 280000000, 75000000, 86000000, 58],
    ["Nomina y Compensaciones", 1200000000, 310000000, 315000000, 52],
    ["Operaciones y Soporte", 180000000, 40000000, 45000000, 47],
    ["Legal y Compliance", 60000000, 15000000, 16000000, 52],
    ["Capacitacion Interna", 40000000, 10000000, 12000000, 55],
    ["Imprevistos", 100000000, 0, 0, 0],
]
for row in data:
    ws.append(row)
total = [sum(r[i] for r in data) for i in range(1, 4)]
ws.append(["TOTAL", total[0], total[1], total[2], ""])
wb.save(str(BASE / "drive" / "fin" / "FIN-002.xlsx"))
print("[creado] FIN-002.xlsx")

# ── COM-001.xlsx ──
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Catalogo"
ws.append(["Codigo", "Curso", "Categoria", "Duracion (h)", "Precio", "Modalidad"])
cursos = [
    ["CUR-001", "Python desde Cero", "Programacion", 40, 149000, "Virtual"],
    ["CUR-002", "JavaScript Avanzado", "Programacion", 30, 129000, "Virtual"],
    ["CUR-003", "SQL y Bases de Datos", "Datos", 25, 99000, "Virtual"],
    ["CUR-004", "Machine Learning Intro", "Datos", 35, 179000, "Virtual"],
    ["CUR-005", "Ingles B1 Intermedio", "Idiomas", 60, 89000, "Mixto"],
    ["CUR-006", "Ingles B2 Avanzado", "Idiomas", 60, 109000, "Mixto"],
    ["CUR-007", "Frances A1-A2", "Idiomas", 50, 89000, "Virtual"],
    ["CUR-008", "Marketing Digital", "Negocios", 20, 79000, "Virtual"],
    ["CUR-009", "Finanzas para no Financieros", "Negocios", 15, 69000, "Virtual"],
    ["CUR-010", "Excel Avanzado", "Ofimatica", 20, 59000, "Virtual"],
    ["CUR-011", "Power BI Dashboard", "Datos", 25, 119000, "Virtual"],
    ["CUR-012", "Desarrollo Web Full Stack", "Programacion", 120, 349000, "Virtual"],
    ["CUR-013", "AWS Cloud Practitioner", "Cloud", 30, 199000, "Virtual"],
    ["CUR-014", "Ingles A1 Principiante", "Idiomas", 50, 69000, "Virtual"],
    ["CUR-015", "Liderazgo y Gestion de Equipos", "Habilidades Blandas", 12, 59000, "Virtual"],
]
for c in cursos:
    ws.append(c)
wb.save(str(BASE / "drive" / "com" / "COM-001.xlsx"))
print("[creado] COM-001.xlsx")

# ── OPE-002.pdf ──
pdf_path = BASE / "drive" / "ope" / "OPE-002.pdf"
doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
styles = getSampleStyleSheet()
story = []
story.append(Paragraph("Plan de Continuidad del Negocio", styles["Title"]))
story.append(Spacer(1, 12))
story.append(Paragraph("Alcance", styles["Heading2"]))
story.append(Paragraph("Este plan aplica a todos los servicios criticos de la plataforma AprendeYa. Su objetivo es garantizar la continuidad operativa ante interrupciones."))
story.append(Paragraph("Escenarios de Contingencia", styles["Heading2"]))
story.append(Paragraph("1. Caida del servidor principal: Activacion automatica de failover a region secundaria (us-west-2). RTO < 5 minutos."))
story.append(Paragraph("2. Ataque DDoS: Mitigacion via AWS Shield Advanced + WAF. RTO < 15 minutos."))
story.append(Paragraph("3. Perdida de datos: Restauracion desde snapshot con retencion de 30 dias. RTO < 2 horas."))
story.append(Paragraph("4. Desastre natural en datacenter: Activacion del DRP en region alterna. RTO < 4 horas."))
story.append(Paragraph("Roles y Responsabilidades", styles["Heading2"]))
story.append(Paragraph("Director de Tecnologia: Activar el plan y coordinar la respuesta. Lider de Infraestructura: Ejecutar la recuperacion tecnica. Lider de Comunicaciones: Informar a afectados."))
story.append(Paragraph("Procedimiento de Activacion", styles["Heading2"]))
story.append(Paragraph("1. Deteccion del incidente. 2. Clasificacion (critico/alto/medio/bajo). 3. Si es critico o alto, el Director activa el plan. 4. Ejecucion de recuperacion segun escenario. 5. Verificacion y retorno a operacion normal. 6. Documentacion de lecciones aprendidas."))
doc.build(story)
print("[creado] OPE-002.pdf")

# ── COM-002.pptx ──
prs = pptx.Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

def set_font(paragraph, size, color=None, bold=False):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    if color:
        paragraph.font.color.rgb = RGBColor(*color)

# Slide 1: Portada
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
txBox = add_textbox(slide, 1, 2.5, 11, 2)
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Estrategia de Marketing Digital 2026"
set_font(p, 40, (255, 255, 255), True)
p.alignment = pptx.enum.text.PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "AprendeYa — Plan Anual"
set_font(p2, 24, (204, 204, 204))
p2.alignment = pptx.enum.text.PP_ALIGN.CENTER

# Slide 2: Objetivos
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = add_textbox(slide, 0.5, 0.5, 12, 6)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Objetivos 2026"
set_font(p, 32, bold=True)
for obj in [
    "Incrementar la base de usuarios activos en un 40% vs 2025",
    "Alcanzar 15,000 estudiantes pagos mensuales",
    "Reducir el CAC en un 20%",
    "Aumentar el NPS de 45 a 65",
    "Expandir presencia en Colombia, Mexico y Peru"
]:
    p2 = tf.add_paragraph()
    p2.text = f"- {obj}"
    set_font(p2, 18)
    p2.space_before = Pt(8)

# Slide 3: Canales
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = add_textbox(slide, 0.5, 0.5, 12, 6)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Canales y Presupuesto"
set_font(p, 32, bold=True)
for ch in [
    "Google Ads (Search + Display): $40,000 USD/mes",
    "Meta Ads (Facebook + Instagram): $30,000 USD/mes",
    "TikTok Ads: $15,000 USD/mes",
    "LinkedIn Ads (cursos corporativos): $10,000 USD/mes",
    "Email Marketing (Mailchimp): $5,000 USD/mes",
    "SEO y Contenido (blog + YouTube): $8,000 USD/mes"
]:
    p2 = tf.add_paragraph()
    p2.text = f"- {ch}"
    set_font(p2, 16)
    p2.space_before = Pt(6)

# Slide 4: KPIs
slide = prs.slides.add_slide(prs.slide_layouts[6])
txBox = add_textbox(slide, 0.5, 0.5, 12, 6)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "KPIs Clave"
set_font(p, 32, bold=True)
for kpi in [
    "CAC: $25 USD actual -> $20 USD meta",
    "LTV: $150 USD actual -> $200 USD meta",
    "Tasa de Conversion: 3.2% -> 4.5% meta",
    "Tasa de Retencion Mensual: 82% -> 88% meta",
    "Visitantes Mensuales: 120K -> 200K meta",
    "Leads Calificados: 3,000/mes -> 5,000/mes meta"
]:
    p2 = tf.add_paragraph()
    p2.text = f"- {kpi}"
    set_font(p2, 16)
    p2.space_before = Pt(6)

prs.save(str(BASE / "sharepoint" / "com" / "COM-002.pptx"))
print("[creado] COM-002.pptx")

print("\n--- TODOS GENERADOS ---")
